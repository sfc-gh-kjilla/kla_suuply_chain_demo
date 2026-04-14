from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

BASE = os.path.dirname(os.path.abspath(__file__))
SS = os.path.join(BASE, "screenshots")
OUT = os.path.join(BASE, "KLA_Escalation_Command_Center_Demo.pptx")

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

DARK_BG = RGBColor(0x0F, 0x17, 0x2A)
ACCENT_BLUE = RGBColor(0x29, 0xB6, 0xF6)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xCC, 0xCC, 0xCC)
SNOWFLAKE_BLUE = RGBColor(0x29, 0xB6, 0xF6)


def set_slide_bg(slide, color=DARK_BG):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_text_box(slide, left, top, width, height, text, font_size=18, color=WHITE, bold=False, alignment=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.alignment = alignment
    return txBox


def add_image_centered(slide, img_path, max_w=12.5, max_h=6.0, top_offset=1.2):
    from PIL import Image
    img = Image.open(img_path)
    w, h = img.size
    aspect = w / h
    target_w = max_w
    target_h = target_w / aspect
    if target_h > max_h:
        target_h = max_h
        target_w = target_h * aspect
    left = (13.333 - target_w) / 2
    slide.shapes.add_picture(img_path, Inches(left), Inches(top_offset), Inches(target_w), Inches(target_h))


def add_notes(slide, text):
    notes_slide = slide.notes_slide
    notes_slide.notes_text_frame.text = text


# ── SLIDE 1: Title ──
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_text_box(slide, 1.5, 1.5, 10, 1.5, "KLA Escalation Command Center", 44, ACCENT_BLUE, True, PP_ALIGN.CENTER)
add_text_box(slide, 1.5, 3.2, 10, 1, "Snowflake + Cortex AI | Real-Time Field Service Intelligence", 24, WHITE, False, PP_ALIGN.CENTER)
add_text_box(slide, 1.5, 4.5, 10, 0.8, "Powered by Snowflake Data Cloud  •  Cortex AI Agent  •  SAP S/4 HANA Integration", 16, LIGHT_GRAY, False, PP_ALIGN.CENTER)
add_text_box(slide, 1.5, 6.0, 10, 0.5, "Prepared for Dave Kelly & Jim Cordova  |  March 2026", 14, LIGHT_GRAY, False, PP_ALIGN.CENTER)
add_notes(slide, """TITLE SLIDE
- Welcome Dave and Jim
- "Today we're going to show you a working prototype of the KLA Escalation Command Center"
- "This is built entirely on Snowflake's platform — Cortex AI for intelligence, with real-time SAP S/4 HANA integration"
- "The goal: replace the email/SharePoint/phone-based case management with a single intelligent command center"
""")


# ── SLIDE 2: Agenda ──
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_text_box(slide, 1.5, 0.8, 10, 1, "Agenda", 36, ACCENT_BLUE, True, PP_ALIGN.CENTER)
agenda_items = [
    "1.  Snowflake Platform Overview — Why Snowflake for KLA",
    "2.  Real-Time AI Workflow Demo — Live Escalation Management",
    "3.  SAP Integration & Interoperability — ERP Data Flows",
    "4.  Next Steps & Discussion"
]
for i, item in enumerate(agenda_items):
    add_text_box(slide, 2.5, 2.2 + i * 1.0, 8, 0.8, item, 22, WHITE, False)
add_notes(slide, """AGENDA
- "We'll start with a quick overview of why Snowflake is the right platform for this"
- "Then we'll do a live demo walking through real escalation scenarios"
- "We'll show the SAP integration — how data flows between S/4 HANA and the command center"
- "And finish with next steps and discussion"
- Keep this slide brief — 30 seconds max
""")


# ── Slide definitions: (screenshot_file, title, notes) ──
slides_data = [
    # 3: Dashboard Overview
    ("01_dashboard_overview.png", "Dashboard Overview — Single Pane of Glass",
     """DASHBOARD OVERVIEW
- "This is the KLA Escalation Command Center — the single pane of glass that replaces email, SharePoint, and phone calls."
- Point out key elements:
  • KPI bar at top: Revenue at Risk ($12.4M), SLA Penalty Exposure ($890K), Avg Case Age (5.2 days), Parts Cost MTD ($234K)
  • Workflow strip showing the 6-step escalation lifecycle
  • Fleet overview (top-left): 5 customers, 13 fabs, all scanners tracked
  • Sensor health (top-right): Real-time scanner diagnostics
  • Case detail (center): Currently selected escalation case
  • Escalation cases (right): Queue sorted by severity
  • Chat panel (far right): AI agent with suggested prompts
- "Everything you see is powered by Snowflake — unified data, real-time queries, Cortex AI intelligence"
"""),

    # 4: KPI Bar
    ("02_kpi_bar.png", "Executive KPI Bar — Real-Time Business Impact",
     """KPI BAR
- "These KPIs are calculated in real-time from your case data, parts costs, and SLA agreements"
- Revenue at Risk: $12.4M — "Total production revenue exposed by active escalation cases"
- SLA Penalty Exposure: $890K — "Contractual penalties accumulating from breached SLAs"
- Avg Case Age: 5.2 days — "How long cases stay open — target is under 3 days for SEV1"
- Parts Cost MTD: $234K — "Month-to-date spend on replacement parts with landed cost awareness"
- Cases/Engineer: 1.4 — "Workload balance across your 4,000 field engineers"
- Batch B-2024-X indicator — "Fleet-wide alert: this batch is the root cause in 5 of 7 cases"
- "Today this data lives in 5 different systems. Here it's one glance."
"""),

    # 5: Workflow Strip
    ("03_workflow_strip.png", "Escalation Workflow — 6-Step Lifecycle",
     """WORKFLOW STRIP
- "Every escalation follows this 6-step lifecycle"
- Issue Detected → Case Created → Parts Identified → Landed Cost Calc → Shipment Executed → Case Resolved
- "The active step is highlighted. For the Samsung case, we're at 'Shipment Executed' — part is in transit"
- "This gives leadership instant visibility into where every case stands in the pipeline"
- "In production, each step transition would be automated with Snowflake Tasks and Streams"
"""),

    # 6: Escalation Cases
    ("04_escalation_cases.png", "Escalation Queue — Priority-Sorted Cases",
     """ESCALATION CASES PANEL
- "Your escalation queue — every open case sorted by severity"
- Each card shows: Case ID, title, customer, fab site, assigned engineer, SLA status
- Color coding: Red = SEV1 (production down), Amber = SEV2 (degraded), Blue = SEV3 (monitoring)
- "Click any case to load its full details, sensor health, and communication timeline"
- SEV1 cases at top: Samsung Pyeongtaek and Renesas Naka — both production impacting
"""),

    # 7: SEV1 Filter
    ("05_sev1_filter.png", "SEV1 Filter — Critical Cases Only",
     """SEV1 FILTER
- "Click SEV1 to isolate just the critical cases — production lines down"
- Only 2 cases visible: Samsung Pyeongtaek (ESC-2026-4281) and Renesas Naka (ESC-2026-4305)
- "This is the view a regional VP uses in their morning standup"
- "Both cases share the same root cause: batch B-2024-X NLO Crystal degradation"
- "Pattern detection like this is impossible when cases live in separate email threads"
"""),

    # 8: Case Detail
    ("06_case_detail.png", "Case Deep-Dive — ESC-2026-4281 Samsung Pyeongtaek",
     """CASE DETAIL
- "Let's drill into the Samsung case"
- Case header: ESC-2026-4281, SEV1, SLA BREACHED (72h)
- Title: "Critical laser power degradation — production impact"
- Key metadata: Samsung Pyeongtaek P3, Scanner SCN-KR-001, Engineer Park Jin-soo (on-site)
- Laser power at 82.3mW — below the 90mW critical threshold
- Two production lines halted, customer requesting ETA
- Communication Timeline: Every interaction in chronological order
  • Field engineer notes, customer emails, escalation events, parts requests
  • "This replaces searching through Outlook to find what was said"
- Action buttons: "AI Case Summary" and "Source Parts"
"""),

    # 9: Sensor Health
    ("07_sensor_health.png", "Sensor Health — Real-Time Scanner Diagnostics",
     """SENSOR HEALTH
- "When we clicked the Samsung case, sensor health auto-loaded for SCN-KR-001"
- 193nm Laser Power: 82.3mW — CRITICAL (below 90mW threshold)
- 8 sensor dimensions monitored: laser power, chamber temp, stage vibration, pressure, beam current, high voltage, coolant flow, optics contamination
- Green = healthy, Yellow = warning, Red = critical
- "Click any metric and the AI agent analyzes readings and correlates to open cases"
- "2-3 metrics flagging simultaneously tells the engineer: crystal problem, not calibration"
"""),

    # 10: Telemetry
    ("08_telemetry.png", "Telemetry — Laser Power Degradation Trend",
     """TELEMETRY TAB
- "Laser power time series for SCN-KR-001"
- Shows degradation trend: power declining from ~105mW to 82mW over 12 days
- Red anomaly markers flag readings below 90mW threshold
- "In production, Cortex AI would detect this trend automatically and trigger a proactive case BEFORE the customer calls"
- "This is predictive maintenance — not reactive firefighting"
"""),

    # 11: Maintenance
    ("09_maintenance.png", "Maintenance History — Pattern Detection",
     """MAINTENANCE HISTORY
- "Every maintenance event for this scanner — corrective and preventive"
- SCN-KR-001 had a corrective repair in August: same issue, same part, same batch B-2024-X
- 8.5 hours of downtime from that previous repair
- "The pattern is clear: B-2024-X crystals failing prematurely across the fleet"
- "Without centralization, this correlation takes weeks. Here, the AI flagged it in seconds"
"""),

    # 12: Source Parts
    ("10_source_parts.png", "Parts Sourcing — Landed Cost Comparison",
     """SOURCE PARTS TAB
- "Where supply chain meets field service"
- Need to ship 994-023 NLO Crystal Assembly to Samsung Pyeongtaek, South Korea
- System compares options by total landed cost including tariffs, duties, VAT
- Singapore Hub: $51,480 landed — 1 day delivery, $0 import duty (ASEAN-Korea FTA)
- Tucson Hub: $58,410 landed — 2 day delivery, $3,600 import duty (8% US-Korea tariff)
- "Singapore saves $7,000 per unit AND arrives a day faster"
- "For a $45,000 crystal, tariff awareness drives the sourcing decision"
- System auto-excludes warehouses with batch B-2024-X stock
- NOTE: For this case path, the display may show batch exclusion message — the shipping cards appear via the chat-triggered sourcing flow
"""),

    # 13: Transfer
    ("11_transfer.png", "Parts Transfer — Inter-Warehouse Movement",
     """PARTS TRANSFER TAB
- "Sometimes the part isn't at the right warehouse"
- Select source warehouse, destination, and quantity
- Example: Moving 1 unit from Tucson to San Jose to cover backordered cases
- "The agent confirms transfer with before/after inventory levels"
- "Identifies which open cases this transfer supports"
"""),

    # 14: Inventory
    ("12_inventory.png", "Global Inventory — Real-Time Stock Visibility",
     """INVENTORY TAB
- "Global inventory status for part 994-023 across all four KLA hubs"
- Tucson Hub: 3 on hand, 1 in transit
- San Jose HQ: 0 on hand, 2 in transit, 4 BACKORDERED — "Red OUT OF STOCK warning"
- Singapore Hub: 2 on hand
- Dresden Hub: 1 on hand, 1 in transit
- "6 units globally, 4 backordered at San Jose"
- "With 5+ cases needing this part, inventory is tight"
- "Click any warehouse card → AI tells you which cases are committed to that stock"
"""),

    # 15: AI Summary
    ("13_ai_summary.png", "AI Summary — Priority Recommendations",
     """AI SUMMARY
- "Click 'AI Summary' at the bottom of the cases panel"
- Agent returns: severity counts, SLA status, priority-ranked recommendations
- CRITICAL: Identifies that 5 of 7 cases share root cause: batch B-2024-X crystals
- "This pattern recognition is impossible when cases live in separate email threads"
- "This replaces your Monday morning case review meeting"
- Root cause analysis, fleet impact assessment, and recommended actions — all automated
"""),

    # 16: Chat with Tool Steps
    ("14_chat_tool_steps.png", "AI Agent — Multi-Tool Orchestration",
     """CHAT PANEL — TOOL ORCHESTRATION
- "The AI agent — powered by Snowflake Cortex — orchestrates multiple tools automatically"
- When you ask a question, watch the tool steps appear:
  • scanner_analyst: queries fleet data, case history, sensor readings
  • tech_docs_search: searches technical documentation and service bulletins
  • ship_part: checks inventory and calculates landed costs
- "Each tool step shows status: running (spinner), completed (checkmark)"
- "The agent decides which tools to use based on your question — you don't need to know the underlying data model"
- "This is the power of Cortex Agent: natural language in, structured intelligence out"
"""),

    # 17: Chat Response
    ("14b_chat_response.png", "AI Agent — Structured Response",
     """CHAT RESPONSE
- "After tool orchestration completes, the agent returns a structured response"
- Includes: case summary, sensor analysis, parts recommendations with costs
- "Notice the response references specific data: scanner IDs, power readings, batch numbers, costs"
- "This isn't a generic chatbot — it's querying your actual operational data in real-time"
- "You can ask follow-up questions and the agent maintains context"
"""),

    # 18: SAP Integration Panel
    ("15_sap_integration_panel.png", "SAP Integration — Real-Time Data Flow",
     """SAP INTEGRATION PANEL
- "This panel shows how data flows between SAP S/4 HANA and the Snowflake command center"
- Three data tables visible:
  • Warehouse inventory data from SAP MM
  • Part recommendations with landed cost calculations
  • Tariff/duty rates by trade corridor
- "SAP provides the transactional backbone: inventory, purchase orders, shipping"
- "Snowflake provides the intelligence layer: AI analysis, pattern detection, cost optimization"
- "Data flows in both directions — read from SAP, write back shipment orders"
- "No rip-and-replace of your ERP — this augments SAP with AI intelligence"
"""),

    # 19: Architecture Overlay
    ("16_architecture_overlay.png", "Architecture — SAP + Snowflake Integration",
     """ARCHITECTURE OVERLAY
- "This is the technical architecture behind what you just saw"
- LEFT: SAP S/4 HANA — source of truth for inventory, orders, shipping, master data
- CENTER: Snowflake Data Cloud — unified data platform
  • Data ingestion via Snowpipe Streaming
  • Cortex AI Layer: LLM inference, search, agent orchestration
  • Real-time analytics and dashboards
- RIGHT: Value propositions — what the business gets
  • Single pane of glass for 3,500+ cases
  • AI-powered pattern detection (batch B-2024-X)
  • Landed cost optimization ($7K savings per shipment)
  • Proactive maintenance vs reactive firefighting
- "Key point: SAP stays as your ERP. Snowflake adds the intelligence layer on top"
- Data flows: SAP → Snowflake (inventory, orders) and Snowflake → SAP (shipment write-backs)
"""),

    # 20: Shipment Modal
    ("17_shipment_modal.png", "Execute Shipment — Order Creation",
     """SHIPMENT MODAL
- "Click 'Execute Shipment' on the recommended source"
- Modal shows: part details, source warehouse, destination, cost breakdown
- "One click creates the shipment order, updates case status, notifies field engineer"
- "The write-back goes to SAP: creates a delivery document in S/4 HANA"
- "No manual data entry, no copy-paste between systems"
"""),

    # 21: Write-back
    ("18_shipment_writeback.png", "SAP Write-Back — Automated Order Processing",
     """SHIPMENT WRITE-BACK
- "After execution, the system confirms: order ID, tracking number, ETA"
- "Simultaneously writes back to SAP S/4 HANA: delivery document, goods movement, billing"
- "The escalation case is automatically updated with shipment status"
- "Field engineer Park Jin-soo gets notified: part arriving tomorrow from Singapore"
- "This closes the loop: from AI recommendation to SAP execution in one workflow"
"""),

    # 22: Fleet Overview
    ("19_fleet_overview.png", "Fleet Overview — Global Scanner Status",
     """FLEET OVERVIEW
- "Your global fleet at a glance: 5 customers, 13 fabs, every scanner tracked"
- TSMC: 4 fabs, 30 scanners, 3 warnings
- Samsung: 3 fabs, 20 scanners, 1 critical
- Intel: 3 fabs, 16 scanners, 1 warning
- SK Hynix: 2 fabs, 10 scanners, 1 warning
- Renesas: 1 fab, 3 scanners, 1 critical
- "Red = critical, Yellow = warning. Click any fab → AI pulls up cases for that site"
"""),

    # 23: Light Theme
    ("20_light_theme.png", "Theme Toggle — Presentation Mode",
     """LIGHT THEME
- "Click the sun/moon icon to toggle themes"
- "Light mode for presentations and executive reviews"
- "Dark mode for the ops center at 2am"
- "Both fully supported with consistent styling"
"""),

    # 24: Customer Filter
    ("21_customer_filter.png", "Customer Filter — Contextual Dashboard",
     """CUSTOMER FILTER
- "Filter the entire dashboard by customer: TSMC, Samsung, Intel, SK Hynix, Renesas"
- "Fleet overview, cases, and metrics all contextualize to the selected customer"
- "Perfect for preparing for a customer quarterly review"
- "Or view all customers for the global operations picture"
"""),
]

# ── Generate content slides ──
for img_file, title, notes in slides_data:
    img_path = os.path.join(SS, img_file)
    if not os.path.exists(img_path):
        print(f"WARNING: Missing {img_file}, skipping")
        continue

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_text_box(slide, 0.4, 0.15, 12.5, 0.8, title, 24, ACCENT_BLUE, True, PP_ALIGN.LEFT)
    add_image_centered(slide, img_path, max_w=12.5, max_h=6.0, top_offset=1.1)
    add_notes(slide, notes)


# ── CLOSING SLIDE ──
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_text_box(slide, 1.5, 1.0, 10, 1, "Key Takeaways", 36, ACCENT_BLUE, True, PP_ALIGN.CENTER)

takeaways = [
    ("Today:", "3,500 cases across email, SharePoint, phone. No single view. No SLA tracking. No pattern detection."),
    ("With This:", "Every case, communication, part, cost calculation — in one place. AI surfaces patterns in seconds, not weeks."),
    ("Landed Cost:", "$7,000 savings per shipment via tariff-aware sourcing. ASEAN-Korea FTA applied automatically."),
    ("Built on Snowflake:", "Unified data platform. Cortex AI Agent. SAP S/4 HANA integration. No new infrastructure."),
]

for i, (label, desc) in enumerate(takeaways):
    add_text_box(slide, 2.0, 2.2 + i * 1.1, 1.5, 0.5, label, 18, ACCENT_BLUE, True)
    add_text_box(slide, 3.5, 2.2 + i * 1.1, 8, 0.5, desc, 16, WHITE)

add_text_box(slide, 1.5, 6.5, 10, 0.5, "Questions?", 28, WHITE, True, PP_ALIGN.CENTER)

add_notes(slide, """CLOSING
- "Let me bring it back to the business problem"
- TODAY: 3,500 cases managed across email, SharePoint, and phone calls. No single view, no SLA tracking, no pattern detection.
- WITH THIS: Every case, every communication, every part, every cost calculation in one place. The AI agent surfaces patterns like batch B-2024-X affecting 5 cases — something that would take weeks to discover manually.
- LANDED COST: When you're shipping a $45,000 crystal, knowing Singapore saves $7,000 over Tucson because of ASEAN-Korea FTA isn't a nice-to-have — it's margin.
- BUILT ON SNOWFLAKE: All data in one platform. Cortex Agent for natural language queries. No new infrastructure. Connects to your SAP S/4 HANA and existing data sources.
- "Questions?"
""")

prs.save(OUT)
print(f"Presentation saved to: {OUT}")
print(f"Total slides: {len(prs.slides)}")
