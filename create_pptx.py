import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.enum.shapes import MSO_SHAPE
from lxml import etree

DK1       = RGBColor(0x26, 0x26, 0x26)
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
DK2       = RGBColor(0x11, 0x56, 0x7F)
SF_BLUE   = RGBColor(0x29, 0xB5, 0xE8)
TEAL      = RGBColor(0x71, 0xD3, 0xDC)
ORANGE    = RGBColor(0xFF, 0x9F, 0x36)
VIOLET    = RGBColor(0x7D, 0x44, 0xCF)
PINK      = RGBColor(0xD4, 0x5B, 0x90)
BODY_GREY = RGBColor(0x5B, 0x5B, 0x5B)
TBL_GREY  = RGBColor(0x71, 0x71, 0x71)
LIGHT_BG  = RGBColor(0xF5, 0xF5, 0xF5)
BORDER    = RGBColor(0xC8, 0xC8, 0xC8)
MUTED     = RGBColor(0xBF, 0xBF, 0xBF)

TEMPLATE = os.path.expanduser("~/templates/snowflake_template.pptx")
SCREENSHOTS_DIR = os.path.join(os.path.dirname(__file__), "screenshots")


def set_ph(slide, idx, text):
    ph = slide.placeholders[idx]
    ph.text = text
    ph.text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    ns = 'http://schemas.openxmlformats.org/drawingml/2006/main'
    bodyPr = ph.text_frame._txBody.find(f'{{{ns}}}bodyPr')
    if bodyPr is None:
        bodyPr = etree.SubElement(ph.text_frame._txBody, f'{{{ns}}}bodyPr')
    t_pos = (ph.top or 0) / 914400
    if t_pos < 0.50:
        bodyPr.set('bIns', '0')
    elif 0.60 < t_pos < 1.20:
        bodyPr.set('tIns', '54864')
    if t_pos < 1.20:
        for para in ph.text_frame.paragraphs:
            pPr = para._p.find(f'{{{ns}}}pPr')
            if pPr is None:
                pPr = etree.SubElement(para._p, f'{{{ns}}}pPr')
                para._p.insert(0, pPr)
            pPr.set('indent', '0')
            pPr.set('marL', '0')


def _pad_body_ph(ph):
    t_pos = (ph.top or 0) / 914400
    if t_pos > 1.20:
        ns = 'http://schemas.openxmlformats.org/drawingml/2006/main'
        bodyPr = ph.text_frame._txBody.find(f'{{{ns}}}bodyPr')
        if bodyPr is None:
            bodyPr = etree.SubElement(ph.text_frame._txBody, f'{{{ns}}}bodyPr')
        bodyPr.set('bIns', '91440')


def set_ph_lines(slide, idx, lines, font_size=None):
    ph = slide.placeholders[idx]
    tf = ph.text_frame
    tf.clear()
    _pad_body_ph(ph)
    lines = [l for l in lines if l.strip()]
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        if font_size:
            p.font.size = Pt(font_size)


def set_ph_sections(slide, idx, sections, heading_size=None, body_size=None):
    ph = slide.placeholders[idx]
    tf = ph.text_frame
    tf.clear()
    _pad_body_ph(ph)
    first = True
    for heading, body_lines in sections:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        p.level = 0
        if not first:
            ns = 'http://schemas.openxmlformats.org/drawingml/2006/main'
            pPr = p._p.find(f'{{{ns}}}pPr')
            if pPr is None:
                pPr = etree.SubElement(p._p, f'{{{ns}}}pPr')
                p._p.insert(0, pPr)
            spcBef = etree.SubElement(pPr, f'{{{ns}}}spcBef')
            spcPts = etree.SubElement(spcBef, f'{{{ns}}}spcPts')
            spcPts.set('val', '1400')
        first = False
        run = p.add_run()
        run.text = heading
        run.font.bold = True
        run.font.color.rgb = DK2
        if heading_size:
            run.font.size = Pt(heading_size)
        for line in body_lines:
            bp = tf.add_paragraph()
            bp.level = 1
            bp.text = line
            if body_size:
                bp.font.size = Pt(body_size)


def add_shape_text(slide, shape_type, left, top, width, height,
                   text, fill_colour, font_colour,
                   font_size=10, bold=False, alignment=PP_ALIGN.CENTER):
    shape = slide.shapes.add_shape(
        shape_type,
        Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_colour
    shape.line.fill.background()
    if width <= 2.0 and '\n' not in text and ' ' in text:
        text = text.replace(' ', '\n')
    tf = shape.text_frame
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Pt(4)
    tf.margin_right = Pt(4)
    tf.margin_top = Pt(2)
    tf.margin_bottom = Pt(2)
    p = tf.paragraphs[0]
    p.text = text
    p.font.name = "Arial"
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = font_colour
    p.alignment = alignment
    return shape


def add_speaker_notes(slide, notes_text):
    notes_slide = slide.notes_slide
    notes_tf = notes_slide.notes_text_frame
    notes_tf.text = notes_text


def add_screenshot_slide(prs, layout_idx, title, subtitle, screenshot_path, notes):
    slide = prs.slides.add_slide(prs.slide_layouts[layout_idx])
    set_ph(slide, 0, title)
    set_ph(slide, 1, subtitle)

    img_top = Inches(1.30)
    img_bottom_max = Inches(5.05)
    max_width = Inches(9.10)
    max_height = img_bottom_max - img_top

    from PIL import Image
    img = Image.open(screenshot_path)
    img_w, img_h = img.size
    aspect = img_w / img_h

    calc_w = max_width
    calc_h = int(calc_w / aspect)
    if calc_h > max_height:
        calc_h = max_height
        calc_w = int(calc_h * aspect)

    left = Inches(0.40) + (max_width - calc_w) // 2
    slide.shapes.add_picture(screenshot_path, left, img_top, calc_w, calc_h)

    add_speaker_notes(slide, notes)
    return slide


prs = Presentation(TEMPLATE)

while len(prs.slides) > 0:
    sldId = prs.slides._sldIdLst[0]
    rId = (sldId.get('{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id')
           or sldId.get('r:id'))
    if rId:
        prs.part.drop_rel(rId)
    prs.slides._sldIdLst.remove(sldId)

print(f"Template loaded, {len(prs.slides)} slides (cleared)")

# ============================================================
# SLIDE 1: Cover
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[13])
set_ph(slide, 3, "KLA ESCALATION\nCOMMAND CENTER")
set_ph(slide, 0, "Powered by Snowflake Cortex AI")
set_ph(slide, 2, "KLA Corporation  |  April 2026")
add_speaker_notes(slide, """Welcome to the KLA Escalation Command Center demo.

This platform demonstrates how KLA can leverage Snowflake's Cortex AI to transform field service operations for their global DUV and EUV semiconductor scanner fleet.

Value proposition: AI-powered diagnostics that reduce mean time to resolution, optimize parts sourcing across global warehouses, and ensure trade compliance — all from a single pane of glass.

The demo covers a real-world scenario: a 193nm DUV laser power degradation detected at Samsung Pyeongtaek, and how the system autonomously diagnoses, sources parts, and manages the escalation through resolution.""")
print("Slide 1: Cover")

# ============================================================
# SLIDE 2: Agenda / Chapter Divider
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[18])
set_ph(slide, 1, "DEMO\nWALKTHROUGH")
add_speaker_notes(slide, """We'll walk through the complete escalation workflow:

1. Dashboard Overview — KPIs, fleet health, SLA breach alerts, and escalation cases
2. Case Management — Detailed case drill-down with AI-powered diagnostics
3. Sensor Health — Multi-metric health view per scanner
4. Telemetry — Real-time time-series data with warning and critical thresholds
5. Maintenance History — Service logs and prior maintenance records
6. Parts Sourcing — Shipping comparison, inventory, transfers, and multi-source landed cost
7. Trade Compliance — Automated 5-step regulatory screening
8. AI Intelligence — Cortex Search and Cortex Agent for natural language operations
9. Customer Filtering — Per-customer views for TSMC, Samsung, and more
10. Architecture — How it's built end-to-end on Snowflake with SAP integration""")
print("Slide 2: Chapter divider")

# ============================================================
# SLIDE 3: Dashboard Overview
# ============================================================
add_screenshot_slide(prs, 0,
    "ESCALATION COMMAND CENTER",
    "Executive dashboard with real-time KPIs, fleet health, and SLA breach alerts",
    os.path.join(SCREENSHOTS_DIR, "01_dashboard_overview.png"),
    """This is the main Escalation Command Center — the single pane of glass for KLA field service operations.

TOP: SLA Breach Alert banner showing 4 cases that have exceeded their SLA window — each one clickable for immediate drill-down.

KPI BAR: Revenue at risk ($5.8M), SLA penalty exposure ($38K), average resolution time (4.4 days), parts cost MTD ($257K), and batch B-2024-X exposure (5 of 7 scanners).

LEFT: Customer Fleet Overview showing TSMC (30 scanners), Samsung (20), Intel (16), SK Hynix, and Renesas with warning/critical counts per customer.

RIGHT: Escalation Cases panel with 7 open cases, filterable by severity (ALL/SEV1/SEV2/SEV3) plus an AI Summary option.

BOTTOM: Tabbed workflow panel with 11 views — Case detail, Sensors, Telemetry, History, Source, Transfer, Inventory, Optimize, Multi-Src, Comply, and Search.

The numbered workflow strip at top shows the 10-step resolution process from Detect through Resolve.""")
print("Slide 3: Dashboard overview")

# ============================================================
# SLIDE 4: Case Detail
# ============================================================
add_screenshot_slide(prs, 0,
    "CASE DRILL-DOWN",
    "Detailed escalation case view with timeline and assigned engineer",
    os.path.join(SCREENSHOTS_DIR, "02_case_detail.png"),
    """When you click on an escalation case, the bottom panel shows the full case detail.

This view displays:
- Case ID, severity, status, and SLA remaining time
- Customer, site, and scanner information
- Assigned field engineer with contact details
- Case description and root cause analysis
- Communication timeline showing all interactions

In this demo, ESC-2026-4281 at Samsung Pyeongtaek is the focal case — a SEV1 laser power degradation that has breached its SLA window. The system has already identified crystal batch B-2024-X as the root cause and is recommending immediate crystal replacement.

Key talking point: Field engineers get full case context without switching between email, SharePoint, and CRM systems — everything is centralized in Snowflake.""")
print("Slide 4: Case detail")

# ============================================================
# SLIDE 5: Sensors
# ============================================================
add_screenshot_slide(prs, 0,
    "SENSOR HEALTH PANEL",
    "Multi-metric health view sorted by worst-to-best condition",
    os.path.join(SCREENSHOTS_DIR, "03_sensors.png"),
    """The Sensors tab provides a comprehensive health view for the selected scanner, showing 8 key metrics sorted from worst to best condition:

- Laser Power (CRITICAL — 82.3mW, below 90mW threshold)
- Chamber Pressure (WARNING)
- Temperature, Vibration, Beam Current, High Voltage, Coolant Flow, Optics Contamination (all HEALTHY)

Each metric shows current value, threshold ranges, and status with color-coded indicators.

Clicking on any metric triggers an AI analysis via the Cortex Agent, providing diagnostic context and recommended actions.

This replaced the original scatter plot visualization based on Tim Long's feedback — "it's just one metric" was the concern with the old avg power display. Now field engineers see the full picture of scanner health at a glance.""")
print("Slide 5: Sensors")

# ============================================================
# SLIDE 6: Telemetry
# ============================================================
add_screenshot_slide(prs, 0,
    "REAL-TIME TELEMETRY",
    "Time-series sensor data with warning and critical thresholds",
    os.path.join(SCREENSHOTS_DIR, "04_telemetry.png"),
    """The Telemetry tab shows real-time sensor readings from the scanner fleet, stored in Snowflake's TELEMETRY_READINGS table (6,258 rows of time-series data).

This view displays the laser power trend over time for scanner SCN-KR-001, showing the gradual degradation that triggered the SEV1 alert. Two horizontal threshold lines are visible:
- Warning threshold at 90mW (yellow)
- Critical threshold at 85mW (red)

The chart uses LTTB (Largest Triangle Three Buckets) downsampling for smooth rendering of large datasets.

Key talking point: All telemetry data is ingested into Snowflake in near-real-time. Cortex Analyst can query this data using natural language — no SQL knowledge needed for field engineers.

The data shows the 193nm laser power dropping from normal operating range (95-100mW) down to the current 82.3mW, crossing the warning threshold approximately 4 days ago.""")
print("Slide 6: Telemetry")

# ============================================================
# SLIDE 7: Maintenance History
# ============================================================
add_screenshot_slide(prs, 0,
    "MAINTENANCE HISTORY",
    "Service logs and prior maintenance records from Snowflake",
    os.path.join(SCREENSHOTS_DIR, "05_maintenance_history.png"),
    """The Maintenance History tab pulls from the MAINTENANCE_LOGS table in Snowflake.

This shows the service history for the selected scanner, including previous crystal replacements, calibration events, and preventive maintenance activities.

Key insight: The AI agent uses this historical data to identify patterns — in this case, it recognized that the crystal from batch B-2024-X has a higher failure rate, connecting this case to a broader fleet-wide issue affecting 5 of 7 scanners with that batch.

This is critical for pattern recognition: rather than treating each case in isolation, Snowflake enables cross-fleet analysis that identifies systemic issues early.""")
print("Slide 7: Maintenance history")

# ============================================================
# SLIDE 8: Source / Shipping
# ============================================================
add_screenshot_slide(prs, 0,
    "SHIPPING COMPARISON",
    "Source part options with cost, transit time, and delivery analysis",
    os.path.join(SCREENSHOTS_DIR, "06_source_shipping.png"),
    """The Source tab provides a shipping comparison for the required replacement part — 994-023 NLO Harmonic Crystal Assembly.

It shows available inventory across KLA's global warehouses with:
- Unit cost and shipping cost per option
- Transit time and delivery windows
- Import duties and tariff rates
- Total landed cost comparison

This demonstrates how Snowflake consolidates supply chain data from multiple ERP systems and provides real-time visibility into parts availability and logistics costs — enabling field engineers to make optimal sourcing decisions in minutes rather than hours.""")
print("Slide 8: Shipping source")

# ============================================================
# SLIDE 9: Part Transfer
# ============================================================
add_screenshot_slide(prs, 0,
    "WAREHOUSE TRANSFER",
    "Warehouse-to-warehouse part transfer management",
    os.path.join(SCREENSHOTS_DIR, "07_part_transfer.png"),
    """The Transfer tab manages warehouse-to-warehouse part transfers.

When a part isn't available at the nearest warehouse, the system can initiate an inter-warehouse transfer to optimize delivery time and cost.

This is backed by Snowflake's WAREHOUSE_INVENTORY and SHIPMENT_ORDERS tables, with a stored procedure (SHIP_PART) that handles the actual transfer orchestration and writes back status updates to SAP S/4 HANA.

Key talking point: The stored procedure takes PART_ID, SCANNER_ID, ALERT_ID, WAREHOUSE_LOCATION, and SHIPPING_METHOD as parameters and creates the shipment order directly in Snowflake — no manual SAP transaction needed.""")
print("Slide 9: Part transfer")

# ============================================================
# SLIDE 10: Parts Inventory
# ============================================================
add_screenshot_slide(prs, 0,
    "PARTS INVENTORY",
    "Global inventory visibility across KLA's regional warehouses",
    os.path.join(SCREENSHOTS_DIR, "08_parts_inventory.png"),
    """The Inventory tab shows a comprehensive view of parts inventory across KLA's 4 regional warehouses (Singapore, Dresden, San Jose, Tucson).

PARTS_INVENTORY table (14 rows) tracks crystal assemblies, optics modules, and critical spares with real-time stock levels, reorder points, and lead times.

Key demo point: The system highlights low-stock items and can automatically trigger replenishment orders through SAP S/4 HANA integration. This consolidated view eliminates the need for field engineers to check multiple regional inventory systems.""")
print("Slide 10: Parts inventory")

# ============================================================
# SLIDE 11: Optimization
# ============================================================
add_screenshot_slide(prs, 0,
    "DYNAMIC SCENARIO MODELING",
    "AI-powered optimization for parts sourcing and logistics",
    os.path.join(SCREENSHOTS_DIR, "09_optimization.png"),
    """The Optimize tab provides dynamic scenario modeling for supply chain decisions.

Users can adjust parameters like shipping method, quantity, and urgency to see the impact on cost, delivery time, and SLA compliance.

This is powered by Snowflake's compute engine running optimization algorithms against real-time inventory and logistics data.

Key value: Field engineers can make data-driven decisions in minutes rather than hours of manual analysis. The system evaluates multiple scenarios simultaneously and recommends the optimal sourcing strategy based on current constraints.""")
print("Slide 11: Optimization")

# ============================================================
# SLIDE 12: Multi-Source
# ============================================================
add_screenshot_slide(prs, 0,
    "MULTI-SOURCE ANALYSIS",
    "Landed cost comparison with ECO-aware part substitution",
    os.path.join(SCREENSHOTS_DIR, "10_multi_source.png"),
    """The Multi-Source tab is one of the most powerful features — it provides a landed cost comparison across multiple sourcing options.

HIGHLIGHT: For part 994-023 shipping to Samsung Pyeongtaek:
- Singapore: Best landed cost with ASEAN-Korea FTA benefits, 1-day transit
- Dresden: Moderate cost with EU-Korea FTA
- Tucson/San Jose: Higher landed cost, 2-day transit, no FTA

The system also shows available part substitutions with ECO-aware validation, ensuring substitutes meet engineering change order requirements.

The "Proceed to Compliance" button flows directly into the trade compliance workflow — connecting sourcing decisions to regulatory screening in a single workflow.""")
print("Slide 12: Multi-source")

# ============================================================
# SLIDE 13: Trade Compliance
# ============================================================
add_screenshot_slide(prs, 0,
    "TRADE COMPLIANCE",
    "5-step automated compliance workflow with regulatory screening",
    os.path.join(SCREENSHOTS_DIR, "11_trade_compliance.png"),
    """The Compliance tab implements a full 5-step trade compliance workflow:

1. Entity Screening — Check against restricted entities lists (Huawei, YMTC, Russian Microelectronics in RESTRICTED_ENTITIES table)
2. Region Validation — Verify destination against RESTRICTED_REGIONS (Russia, North Korea, Iran, etc.)
3. Tariff Analysis — Calculate duties using TARIFF_RATES (15 rate entries)
4. Documentation — Generate compliance certificates
5. Approval — Final sign-off with audit trail in COMPLIANCE_LOG

BLOCKED entities trigger a full-screen modal with red border showing entity details and legal references (EAR §744 / OFAC 31 CFR 560). Green checkmarks indicate PASSED steps.

This is critical for KLA's semiconductor equipment exports, which are subject to strict US export controls (EAR/ITAR) and international trade regulations.

Key talking point: Snowflake stores all compliance data and audit trails in one place, enabling real-time screening rather than batch processing.""")
print("Slide 13: Trade compliance")

# ============================================================
# SLIDE 14: Cortex AI Search
# ============================================================
add_screenshot_slide(prs, 0,
    "CORTEX AI SEARCH",
    "RAG-powered technical documentation search for field engineers",
    os.path.join(SCREENSHOTS_DIR, "12_cortex_search.png"),
    """The Search tab demonstrates Snowflake Cortex Search — a fully managed RAG (Retrieval Augmented Generation) service.

TECHNICAL_DOCS_SEARCH is a Cortex Search Service indexing 8 technical documents including maintenance manuals, service bulletins, and troubleshooting guides.

Embedding model: snowflake-arctic-embed-m-v1.5

Field engineers can ask natural language questions like "How do I replace the NLO crystal assembly?" and get precise answers grounded in KLA's actual technical documentation — with source citations.

Key differentiator: Unlike generic AI chatbots, this is grounded in KLA's proprietary technical knowledge base stored securely in Snowflake. No data leaves the Snowflake environment.""")
print("Slide 14: Cortex search")

# ============================================================
# SLIDE 15: AI Chat
# ============================================================
add_screenshot_slide(prs, 0,
    "AI DIAGNOSTIC AGENT",
    "Cortex Agent with multi-tool intelligence for field operations",
    os.path.join(SCREENSHOTS_DIR, "13_ai_chat_response.png"),
    """The AI Chat panel shows the KLA Diagnostic Agent in action — a Cortex Agent with 3 integrated tools:

1. scanner_analyst (Cortex Analyst) — Text-to-SQL queries against the KLA_SCANNER_INTELLIGENCE semantic view
2. tech_docs_search (Cortex Search) — RAG over technical documentation
3. ship_part (Stored Procedure) — Autonomous part shipping and order creation

In this screenshot, the agent has responded to "Summarize open escalation cases by priority" — showing SEV1 cases with SLA breach status and recommended actions.

The agent has 5 suggested prompts covering case summaries, SLA analysis, batch impact, and part availability — all accessible through natural language.

Key demo point: The agent can autonomously diagnose issues, search documentation, query the database, and even initiate part shipments — all through conversation. This is the AI copilot vision for 4,000+ KLA field engineers.""")
print("Slide 15: AI chat")

# ============================================================
# SLIDE 16: SEV1 Filtered
# ============================================================
add_screenshot_slide(prs, 0,
    "SEV1 ESCALATION FILTER",
    "Priority filtering with SLA breach tracking",
    os.path.join(SCREENSHOTS_DIR, "14_sev1_filtered.png"),
    """The escalation cases panel with SEV1 filter applied, showing only the highest-priority cases.

Each case card displays:
- Customer badge (teal pill with customer name)
- Case ID and severity classification
- Assigned field engineer
- SLA remaining time with breach indicator
- Brief description

The panel supports ALL/SEV1/SEV2/SEV3 filtering and shows real-time SLA breach tracking with color-coded urgency indicators. Breached cases show negative hours remaining.

This directly addresses the current pain point Todd Smith (EY) identified: KLA manages ~3,500 active cases via email and SharePoint with no centralized priority view.""")
print("Slide 16: SEV1 filtered")

# ============================================================
# SLIDE 17: Customer Filter — TSMC
# ============================================================
add_screenshot_slide(prs, 0,
    "CUSTOMER VIEW: TSMC",
    "Filtered dashboard showing TSMC-specific fleet and cases",
    os.path.join(SCREENSHOTS_DIR, "15_customer_tsmc.png"),
    """The customer dropdown filter applied to TSMC — KLA's largest customer representing approximately 50% of revenue.

When TSMC is selected:
- Fleet overview shows only TSMC's 30 scanners across their fabs
- Escalation cases filter to TSMC-specific cases only
- KPIs update to reflect TSMC's revenue at risk and SLA exposure
- All downstream tabs (telemetry, inventory, etc.) scope to TSMC context

This was a key request from Jon Cernok (AE): "TSMC is 50% of revenue — they need to be able to see their own view."

Key talking point for Dave Kelly and Jim Cordova: Each customer or regional manager gets a personalized view without building separate dashboards.""")
print("Slide 17: Customer TSMC")

# ============================================================
# SLIDE 18: Customer Filter — Samsung
# ============================================================
add_screenshot_slide(prs, 0,
    "CUSTOMER VIEW: SAMSUNG",
    "Filtered dashboard showing Samsung-specific fleet and escalations",
    os.path.join(SCREENSHOTS_DIR, "16_customer_samsung.png"),
    """Same customer filter applied to Samsung — the focal customer in our demo scenario.

The Samsung view highlights:
- 20 scanners across Samsung fabs (Pyeongtaek, Hwaseong, Austin)
- The critical ESC-2026-4281 case at Pyeongtaek Fab 18 (our main demo case)
- Samsung-specific SLA metrics and breach status

This demonstrates how the same dashboard scales across KLA's customer base — each customer gets their own context without duplicate infrastructure.

The dropdown supports: All Customers (default), TSMC, Samsung, Intel, SK Hynix, and Renesas.""")
print("Slide 18: Customer Samsung")

# ============================================================
# SLIDE 19: Architecture
# ============================================================
add_screenshot_slide(prs, 0,
    "SYSTEM ARCHITECTURE",
    "End-to-end Snowflake architecture powering the platform",
    os.path.join(SCREENSHOTS_DIR, "17_architecture_overlay.png"),
    """The Architecture overlay shows the complete Snowflake-powered system architecture:

DATA LAYER: 16 tables in KLA_SUPPLY_CHAIN.APP schema including SCANNERS, TELEMETRY_READINGS, ALERTS, PARTS_INVENTORY, ESCALATION_DEMANDS, and compliance tables.

AI LAYER:
- Cortex Agent (KLA_DIAGNOSTIC_AGENT) orchestrating 3 tools
- Cortex Analyst with KLA_SCANNER_INTELLIGENCE semantic view for text-to-SQL
- Cortex Search (TECHNICAL_DOCS_SEARCH) for RAG over technical docs
- snowflake-arctic-embed-m-v1.5 embedding model

INTEGRATION: SAP S/4 HANA connection for ERP data synchronization, stored procedures for part shipping and writeback.

FRONTEND: React + TypeScript application consuming Snowflake REST APIs.

Key talking point: Everything runs on Snowflake — no separate ML infrastructure, no external vector databases, no additional API services. One platform for data, AI, and applications.""")
print("Slide 19: Architecture")

# ============================================================
# SLIDE 20: SAP Integration
# ============================================================
add_screenshot_slide(prs, 0,
    "SAP S/4 HANA INTEGRATION",
    "Enterprise ERP integration for supply chain data flow",
    os.path.join(SCREENSHOTS_DIR, "18_sap_integration.png"),
    """The SAP S/4 HANA integration panel shows how the platform connects to KLA's enterprise ERP system.

Data flows between Snowflake and SAP include:
- Material master data synchronization
- Purchase order creation and tracking
- Inventory movements and stock transfers
- Maintenance order creation for field service
- Compliance documentation and audit trails

This demonstrates that Snowflake serves as the intelligent layer on top of existing ERP investments, adding AI capabilities without disrupting current SAP workflows.

Key for CIO audience (Dave Kelly, Jim Cordova): This is additive to SAP, not a replacement. Snowflake enhances SAP data with AI-powered analytics and autonomous operations. No SAP customization required.""")
print("Slide 20: SAP integration")

# ============================================================
# SLIDE 21: Key Capabilities Summary (3-column)
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[7])
set_ph(slide, 0, "KEY CAPABILITIES")
set_ph(slide, 4, "Snowflake Cortex AI powering field service operations")
set_ph_sections(slide, 1, [
    ("AI-Powered Diagnostics", [
        "Cortex Agent with 3 integrated tools",
        "Natural language queries via Cortex Analyst",
        "RAG over technical documentation",
        "Autonomous part shipping via stored procedures",
    ]),
], heading_size=11, body_size=10)
set_ph_sections(slide, 2, [
    ("Supply Chain Intelligence", [
        "Multi-source landed cost analysis",
        "Global inventory across 4 warehouses",
        "ECO-aware part substitution",
        "Dynamic scenario optimization",
    ]),
], heading_size=11, body_size=10)
set_ph_sections(slide, 3, [
    ("Compliance & Governance", [
        "5-step trade compliance workflow",
        "Entity and region screening (EAR/OFAC)",
        "Automated tariff calculation",
        "Full audit trail in Snowflake",
    ]),
], heading_size=11, body_size=10)
add_speaker_notes(slide, """Summary of the three key capability pillars:

1. AI-POWERED DIAGNOSTICS: The Cortex Agent autonomously diagnoses scanner issues, searches technical documentation, queries operational data, and can even ship parts — all through natural language. Built with Cortex Agent, Cortex Analyst (semantic view), and Cortex Search (RAG).

2. SUPPLY CHAIN INTELLIGENCE: Multi-source analysis with landed cost comparison across Singapore, Dresden, San Jose, and Tucson warehouses. ECO-aware substitution ensures replacements meet engineering change order requirements. Dynamic optimization modeling evaluates multiple scenarios simultaneously.

3. COMPLIANCE & GOVERNANCE: Automated 5-step trade compliance workflow covering entity screening, region validation, tariff analysis, documentation, and approval. BLOCKED modal for restricted entities (Huawei, YMTC) with legal references. Critical for semiconductor equipment exports subject to US export controls.

All three pillars run entirely on Snowflake — no external AI infrastructure required.""")
print("Slide 21: Key capabilities")

# ============================================================
# SLIDE 22: Business Impact (stat callouts)
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[0])
set_ph(slide, 0, "BUSINESS IMPACT")
set_ph(slide, 1, "Measurable outcomes from AI-powered operations")

stats = [
    ("60%", "Faster\nResolution", "AI diagnostics reduce\nMTTR from 4.4 to 1.8 days"),
    ("$2.1M", "Annual\nSavings", "Optimized sourcing and\nreduced SLA penalties"),
    ("5x", "Engineer\nProductivity", "Natural language replaces\nmanual analysis"),
    ("100%", "Compliance\nCoverage", "Automated screening\nfor every shipment"),
]

for i, (number, label, desc) in enumerate(stats):
    x = 0.40 + i * 2.35
    add_shape_text(slide, MSO_SHAPE.ROUNDED_RECTANGLE, x, 1.50, 2.10, 1.10,
                   number, SF_BLUE, WHITE, font_size=28, bold=True)
    add_shape_text(slide, MSO_SHAPE.ROUNDED_RECTANGLE, x, 2.70, 2.10, 0.55,
                   label, DK2, WHITE, font_size=10, bold=True)
    add_shape_text(slide, MSO_SHAPE.ROUNDED_RECTANGLE, x, 3.35, 2.10, 0.70,
                   desc, LIGHT_BG, DK1, font_size=9)

add_speaker_notes(slide, """Business impact metrics for C-level presentation:

60% FASTER RESOLUTION: AI diagnostics with Cortex Agent reduce mean time to resolution from 4.4 days to approximately 1.8 days by automating root cause analysis, pattern detection across the fleet, and optimal parts sourcing.

$2.1M ANNUAL SAVINGS: Optimized multi-source parts procurement using landed cost analysis, reduced SLA penalty exposure (currently $38K), and fewer emergency shipments through proactive batch analysis.

5x ENGINEER PRODUCTIVITY: 4,000+ field engineers use natural language instead of manual SQL queries, email chains, and spreadsheet analysis. One question to the Cortex Agent replaces hours of data gathering across SAP, SharePoint, and email.

100% COMPLIANCE COVERAGE: Every parts shipment goes through automated 5-step trade compliance screening — entity screening, region validation, tariff calculation, documentation, and approval. No manual gaps in export control coverage.

These are projected metrics based on the demo scenario and can be customized to KLA's actual operational data.""")
print("Slide 22: Business impact")

# ============================================================
# SLIDE 23: Thank You
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[28])
set_ph(slide, 1, "THANK\nYOU")
add_speaker_notes(slide, """Thank you for attending this demo of the KLA Escalation Command Center.

Next steps:
1. Schedule a deeper technical dive with Snowflake engineering
2. Identify a pilot use case within KLA's field service operations (suggested: SEV1 escalation workflow)
3. Connect to KLA's actual SAP and telemetry data sources for a POC
4. Define success metrics with Dave Kelly and Jim Cordova

Contact: Your Snowflake account team for follow-up discussions.

All demo assets including the application code, Snowflake objects (16 tables, semantic view, Cortex Search service, Cortex Agent), and this presentation are available for reference.""")
print("Slide 23: Thank you")

# ============================================================
# Save
# ============================================================
output_dir = os.path.join(os.path.dirname(__file__), "outputs")
os.makedirs(output_dir, exist_ok=True)
output_path = os.path.join(output_dir, "KLA_Supply_Chain_Demo.pptx")
prs.save(output_path)
print(f"\nSaved: {output_path}")
print(f"Total slides: {len(prs.slides)}")
